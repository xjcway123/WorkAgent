"""
LangGraph v1.0 多智能体系统 —— 起步骨架
功能:总结 / 翻译 / 发邮件 / 做 PPT / 做 Excel

架构:Supervisor(主管)路由到 5 个专家智能体。
依赖:pip install langchain langgraph langgraph-supervisor langchain-openai python-pptx openpyxl

环境变量(用 .env + python-dotenv 或直接 export):
    ARK_API_KEY   火山方舟密钥
    SMTP_HOST / SMTP_USER / SMTP_PASS   发邮件用(可选,不发邮件可不配)
"""

import os
import asyncio

from langchain_openai import ChatOpenAI
from langchain_core.tools import tool
from langchain.agents import create_agent              # ← v1.0 当前写法(取代 create_react_agent)
from langgraph_supervisor import create_supervisor
from pydantic import BaseModel, Field
import smtplib
from email.mime.text import MIMEText
import re
from openpyxl import Workbook
from openpyxl.styles import Font
from skill_loader import SkillLoader
from pathlib import Path

from dotenv import load_dotenv

load_dotenv() 


WORKDIR = Path.cwd()
SKILLS_DIR = WORKDIR / "skills"

SKILL_LOADER = SkillLoader(SKILLS_DIR)

# ----------------------------------------------------------------------
# 1. 模型:指向火山方舟(OpenAI 兼容),model 填你的接入点 ID
# ----------------------------------------------------------------------
model = ChatOpenAI(
    model=os.environ["MODEL_NAME"],
    base_url=os.environ["BASE_URL"],
    api_key=os.environ["ARK_API_KEY"],
    temperature=0,
)

OUTPUT_DIR = "./output"
os.makedirs(OUTPUT_DIR, exist_ok=True)



class EmailInput(BaseModel):
    to: str = Field(description="收件人邮箱地址")
    subject: str = Field(description="邮件主题")
    body: str = Field(description="邮件正文内容")


# ----------------------------------------------------------------------
# 2. 工具(tool):被对应的专家智能体调用。docstring 很重要,LLM 靠它判断怎么用
# ----------------------------------------------------------------------

# ----------------------------------------------------------------------
# 2a. Human-in-the-Loop (HITL) broker: 用于让 agent 在执行敏感操作前暂停并
#     等待用户确认。`request_human_input` 工具是入口;`astream_supervisor`
#     通过共享 asyncio.Queue 把请求事件推给前端,并阻塞当前 tool 协程
#     等待 `/chat/respond` 端点 resolve 对应 Future。
# ----------------------------------------------------------------------
class _HitlBroker:
    def __init__(self) -> None:
        self.queue: asyncio.Queue | None = None
        self.pending: dict[str, asyncio.Future] = {}
        self.loop: asyncio.AbstractEventLoop | None = None

    def attach(self, queue: asyncio.Queue, loop: asyncio.AbstractEventLoop) -> None:
        self.queue = queue
        self.loop = loop
        self.pending.clear()

    def detach(self) -> None:
        # Cancel any still-pending requests (e.g. client disconnected)
        for fut in self.pending.values():
            if not fut.done():
                fut.set_exception(asyncio.CancelledError("hitl session ended"))
        self.pending.clear()
        self.queue = None
        self.loop = None

    @property
    def active(self) -> bool:
        return self.queue is not None and self.loop is not None

    def resolve(self, request_id: str, response: str) -> bool:
        fut = self.pending.pop(request_id, None)
        if fut and not fut.done():
            fut.set_result(response)
            return True
        return False


_hitl_broker = _HitlBroker()


def _supervisor_pre_model_hook(state):
    """Force the supervisor to stop looping on the same agent.

    LangGraph's supervisor pattern re-invokes the supervisor LLM after each
    sub-agent returns. The LLM often mis-decides and calls the same agent
    again, generating duplicate output. This hook runs *before* every
    supervisor LLM call. If the supervisor has already routed once via
    ``transfer_to_*``, we collapse the visible history down to
    ``[user_msg, last_agent_response]`` so the LLM only sees a clean
    "user asked X, expert answered Y" view and either passes Y through
    (and the graph ENDs) or routes to a *different* agent for a multi-
    step task. Either way, the same agent can never be called twice.
    """
    messages = state.get("messages", [])
    if not messages:
        return state

    # Count how many times the supervisor has already routed.
    routed = 0
    for m in messages:
        tc = getattr(m, "tool_calls", None) or []
        for call in tc:
            name = call.get("name") if isinstance(call, dict) else getattr(call, "name", "")
            if name and name.startswith("transfer_to_"):
                routed += 1

    if routed == 0:
        # First time the supervisor LLM runs — let it see the full
        # history so it can make a sensible first routing decision.
        return state

    # Already routed at least once. Find the most recent agent response
    # (an AIMessage with content but no tool calls).
    last_response = None
    for m in reversed(messages):
        cls = m.__class__.__name__
        if cls in ("HumanMessage", "ToolMessage", "SystemMessage"):
            continue
        has_tools = bool(getattr(m, "tool_calls", None))
        has_content = bool(getattr(m, "content", ""))
        if has_content and not has_tools:
            last_response = m
            break

    if last_response is None:
        return state

    # Keep the original user turn + the last agent's response. The
    # supervisor LLM will then either pass it through (END) or pick a
    # different agent for the next step — never re-invoke the same one.
    first_user = next(
        (m for m in messages if m.__class__.__name__ == "HumanMessage"),
        None,
    )
    stripped = [first_user, last_response] if first_user else [last_response]
    return {"messages": stripped}


def _is_handoff_tool(name: str) -> bool:
    """LangGraph's supervisor injects a family of internal handoff tools
    (`transfer_to_<agent>`, `transfer_back_to_<agent>`) to move control
    between the supervisor and sub-agents. They are plumbing, not user-
    relevant actions, so the stream consumer skips them.
    """
    if not name:
        return False
    n = name.lower()
    return n.startswith("transfer_to_") or n.startswith("transfer_back_to_")


@tool
async def request_human_input(question: str, options: list[str] | None = None) -> str:
    """请求用户输入或确认 — Human-in-the-Loop 入口。

    Agent 在执行敏感或不可逆操作前(如发邮件、覆盖文件)必须先调用本工具
    暂停执行,等待用户在 UI 上做出选择。**调用方必须 await 本工具**。

    参数:
        question: 提问或确认提示,支持换行(展示给用户)。
        options: 选项列表,如 ["Approve", "Reject"];为 None/空时表示允许
                 用户输入自由文本(前端会渲染为文本框)。

    返回:
        用户选择的选项(如 "Approve"),或用户输入的文本。
        若 HITL 不可用(没有活跃会话),返回固定字符串 "HitlNotAvailable"。
    """
    if not _hitl_broker.active:
        return "HitlNotAvailable"
    import uuid as _uuid
    request_id = _uuid.uuid4().hex
    fut = _hitl_broker.loop.create_future()
    _hitl_broker.pending[request_id] = fut
    await _hitl_broker.queue.put({
        "event": "human_input_required",
        "request_id": request_id,
        "question": question,
        "options": list(options) if options else [],
    })
    return await fut


@tool
async def send_email(to: str, subject: str, body: str) -> str:
    """发送一封纯文本邮件(高风险操作,自动请求用户确认)。

    实际发送前会通过 `request_human_input` 暂停并向用户确认收件人/主题/正文;
    用户选择 "Approve" 才真正发送,选择 "Reject" 或其它选项则取消。

    参数:
        to: 收件人邮箱地址,例如 someone@example.com
        subject: 邮件主题
        body: 邮件正文内容(纯文本)
    返回:
        发送结果说明(成功提示、用户取消、或失败原因)。
    """
    # Ask for user confirmation
    approval = await request_human_input.ainvoke({
        "question": (
            f"📧 确认发送邮件?\n"
            f"收件人: {to}\n"
            f"主题  : {subject}\n"
            f"\n—— 正文预览 ——\n"
            f"{body[:300]}{'...' if len(body) > 300 else ''}"
        ),
        "options": ["Approve", "Reject"],
    })
    if approval != "Approve":
        return f"邮件已取消(用户选择:{approval})"

    msg = MIMEText(body, "plain", "utf-8")
    msg["Subject"] = subject
    msg["From"] = os.environ["SMTP_USER"]
    msg["To"] = to

    port = int(os.getenv("SMTP_PORT", 465))
    try:
        with smtplib.SMTP_SSL(os.environ["SMTP_HOST"], port) as s:
            s.login(os.environ["SMTP_USER"], os.environ["SMTP_PASS"])
            s.send_message(msg)
        return f"已成功发送邮件给 {to}(主题:{subject})"
    except smtplib.SMTPAuthenticationError:
        return "发送失败:邮箱认证失败,请检查 SMTP_USER 和授权码(SMTP_PASS)是否正确"
    except Exception as e:
        return f"发送失败:{e}"


@tool
def make_excel(filename: str, headers: list[str], rows: list[list]) -> str:
    """生成一个 Excel 文件并返回路径。

    参数:
        filename: 文件名(不含后缀),例如 "sales"
        headers: 表头列表,例如 ["月份", "销售额"]
        rows: 二维列表,每个子列表是一行,顺序对应 headers,
              例如 [["1月", 100], ["2月", 200]]
    返回:
        生成结果说明(成功路径或失败原因)。
    """
    # 文件名清洗:非法/路径字符替换为 _,防止报错或路径穿越
    safe_name = re.sub(r"[^\w\-]", "_", filename)
    path = os.path.join(OUTPUT_DIR, f"{safe_name}.xlsx")

    try:
        wb = Workbook()
        ws = wb.active

        # 写表头并加粗
        ws.append(headers)
        for cell in ws[1]:
            cell.font = Font(bold=True)

        # 写数据,逐行校验列数与表头一致
        for i, row in enumerate(rows, start=1):
            if len(row) != len(headers):
                return f"生成失败:第 {i} 行有 {len(row)} 列,与表头 {len(headers)} 列不符"
            ws.append(row)

        wb.save(path)
        return f"Excel 已生成:{path}"
    except Exception as e:
        return f"生成 Excel 失败:{e}"


@tool
def make_ppt(filename: str, title: str, slides: list[dict]) -> str:
    """生成一个 PPT。title=封面标题, slides=内容页列表,每项形如 {"title": "页标题", "content": "页正文"}。返回文件路径。"""
    from pptx import Presentation

    prs = Presentation()
    # 封面页
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    # 内容页
    for slide in slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = slide.get("title", "")
        s.placeholders[1].text = slide.get("content", "")
    path = os.path.join(OUTPUT_DIR, f"{filename}.pptx")
    prs.save(path)
    return f"PPT 已生成:{path}"


@tool
def make_business_ppt(
    filename: str,
    title: str,
    sections: list[dict],
    subtitle: str = "",
    presenter: str = "",
    date: str = "",
    closing_contact: str = "",
) -> str:
    """生成一份按品牌蓝风格统一规范的商务汇报 PPT(16:9, 6 种页型)。

    参数:
        filename: 文件名(不含后缀),如 "sales_report"
        title: 主标题(封面与页脚复用)
        subtitle: 副标题(可选,如"2025 Q3 季度复盘")
        presenter: 汇报人姓名(可选)
        date: 汇报日期(可选,如"2025-10-15")
        sections: 章节列表,每项形如:
            {
                "name": "业务回顾",
                "name_en": "Business Review",          # 可选,英文副名
                "slides": [                              # 该章节的内容页
                    {
                        "type": "kpi",                   # 关键指标页
                        "title": "核心业务指标",
                        "intro": "本季度核心指标稳健增长",   # 可选引言
                        "kpis": [
                            {"label": "总营收", "value": "¥1,280万", "delta": "+12.4% YoY"},
                            ...
                        ],
                        "bullets": ["关键解读 1", "关键解读 2"]  # 可选,卡片下要点
                    },
                    {
                        "type": "bullets",               # 要点页
                        "title": "重点工作",
                        "bullets": ["要点 1", "要点 2", ...]
                    },
                    {
                        "type": "table",                 # 表格页
                        "title": "区域销售对比",
                        "headers": ["区域", "营收", "同比"],
                        "rows": [
                            ["华东", "¥420万", "+15%"],
                            ...
                        ]
                    }
                ]
            }
        closing_contact: 结束页联系方式(可选)

    返回:生成结果说明(成功路径或失败原因)。

    页型:
        1 封面 → 2 目录 → 3 章节分隔页(章节数 ≥ 4 时) → 4 内容页(按 sections 顺序) → 5 结束页
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    # ------------------------------------------------------------------
    # 调色板与字号(全部来自 references/style-guide.md)
    # ------------------------------------------------------------------
    DEEP_BLUE  = RGBColor(0x29, 0x37, 0x64)   # #293764
    BRAND      = RGBColor(0x4A, 0x70, 0xAE)   # #4A70AE
    ACCENT     = RGBColor(0x6E, 0x9F, 0xCB)   # #6E9FCB
    LIGHT_BG   = RGBColor(0xF2, 0xF6, 0xFC)   # #F2F6FC
    DARK_TEXT  = RGBColor(0x1F, 0x1F, 0x1F)   # #1F1F1F
    GRAY_TEXT  = RGBColor(0x80, 0x80, 0x80)   # #808080
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

    CN_FONT = "微软雅黑"
    EN_FONT = "Arial"

    # ------------------------------------------------------------------
    # 工具函数
    # ------------------------------------------------------------------
    def _set_run_font(run, size: int, bold: bool = False, color: RGBColor = DARK_TEXT):
        run.font.name = EN_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        # 中文走 rPr 之外的 eastAsia 字段
        rPr = run._r.get_or_add_rPr()
        from pptx.oxml.ns import qn
        for tag in ("ea", "cs"):
            existing = rPr.find(qn(f"a:{tag}"))
            if existing is not None:
                rPr.remove(existing)
        ea = rPr.makeelement(qn("a:ea"), {"typeface": CN_FONT})
        rPr.append(ea)

    def _add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=DARK_TEXT,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None, line=None):
        """添加一个干净文本框(占位符 placeholder 也走这里,便于完全控制)。"""
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = anchor
        if fill is not None:
            box.fill.solid()
            box.fill.fore_color.rgb = fill
        if line is None:
            box.line.fill.background()
        else:
            box.line.color.rgb = line
            box.line.width = Pt(0.5)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_run_font(run, size, bold, color)
        return box

    def _add_paragraphs(tf, lines, *, size=14, bold=False, color=DARK_TEXT,
                        align=PP_ALIGN.LEFT, line_spacing=1.15, space_after_pt=4):
        """清空并填充多行段落。lines 可以是 str(单行) 或 list[str]。"""
        tf.clear()
        if isinstance(lines, str):
            lines = [lines]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            p.space_after = Pt(space_after_pt)
            run = p.add_run()
            run.text = line
            _set_run_font(run, size, bold, color)

    def _add_rect(slide, x, y, w, h, fill, line=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(0.5)
        shp.shadow.inherit = False
        return shp

    def _add_footer(slide, page_no, total):
        """所有非封面页通用:汇报主题(左) + 页码(右)"""
        # 顶部细色条(内容页用)
        if page_no > 0 and page_no != total:
            _add_rect(slide, 0, 0, 13.333, 0.18, BRAND)
        # 页脚文字
        _add_text(slide, 0.6, 7.18, 8, 0.3, title, size=9, color=GRAY_TEXT, align=PP_ALIGN.LEFT)
        _add_text(slide, 11.5, 7.18, 1.2, 0.3, f"{page_no} / {total}", size=9,
                  color=GRAY_TEXT, align=PP_ALIGN.RIGHT)

    # ------------------------------------------------------------------
    # 准备 Presentation (16:9)
    # ------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # 空白版式,完全自绘

    safe_name = re.sub(r"[^\w\-]", "_", filename)
    path = os.path.join(OUTPUT_DIR, f"{safe_name}.pptx")

    # ------------------------------------------------------------------
    # 预先计算总页数
    # ------------------------------------------------------------------
    content_slide_count = sum(len(sec.get("slides", [])) for sec in sections)
    insert_dividers = len(sections) >= 4
    # 1 cover + 1 toc + (sections * 1 divider) + content + 1 closing
    total_pages = 1 + 1 + (len(sections) if insert_dividers else 0) + content_slide_count + 1

    page_no = 0  # 当前页号(用于页脚)

    # ------------------------------------------------------------------
    # 页型 1:封面
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    # 左侧 5 英寸深蓝大色块
    _add_rect(slide, 0, 0, 5, 7.5, DEEP_BLUE)
    # 顶部装饰细条
    _add_rect(slide, 5, 0, 0.04, 7.5, BRAND)
    # 装饰点(底部)
    for i, ratio in enumerate([0.15, 0.32, 0.5]):
        _add_rect(slide, 5 + ratio, 7.0, 0.08, 0.08, ACCENT)
    # 主标题(在深蓝块上)
    _add_text(slide, 0.5, 2.6, 4.0, 1.2, title, size=36, bold=True, color=WHITE)
    if subtitle:
        # 副标题以 "—— 副标题" 形式
        _add_text(slide, 0.5, 3.8, 4.0, 0.6, f"—— {subtitle}", size=18, bold=True, color=WHITE)
    # 右侧(白底)放品牌主色大字
    _add_text(slide, 5.5, 2.8, 7.5, 1.4, "Business", size=54, bold=True, color=DEEP_BLUE)
    _add_text(slide, 5.5, 4.0, 7.5, 1.4, "Report", size=54, bold=True, color=BRAND)
    # 汇报人 + 日期(深蓝块上,白字)
    info_line = []
    if presenter: info_line.append(f"汇报人:{presenter}")
    if date: info_line.append(f"日期:{date}")
    if info_line:
        _add_text(slide, 0.5, 6.7, 4.0, 0.4, "  ·  ".join(info_line), size=12, color=WHITE)
    page_no += 1

    # ------------------------------------------------------------------
    # 页型 2:目录
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.333, 0.18, BRAND)
    _add_text(slide, 0.6, 0.45, 12, 0.6, "目  录  /  Contents", size=24, bold=True, color=DEEP_BLUE)

    # 目录条目
    cur_y = 1.6
    for i, sec in enumerate(sections):
        idx = i + 1
        # 编号(品牌蓝,32pt)
        _add_text(slide, 0.8, cur_y, 1.0, 0.6, f"{idx:02d}", size=32, bold=True, color=BRAND)
        # 章节名(深灰,18pt)
        _add_text(slide, 2.0, cur_y + 0.05, 8.0, 0.5, sec.get("name", ""), size=18, bold=False, color=DARK_TEXT)
        # 分割线
        _add_rect(slide, 0.8, cur_y + 0.65, 11.7, 0.012, BRAND)
        cur_y += 0.95

    page_no += 1
    _add_footer(slide, page_no, total_pages)

    # ------------------------------------------------------------------
    # 页型 3 + 4:章节分隔页 + 内容页
    # ------------------------------------------------------------------
    for s_idx, sec in enumerate(sections):
        # 章节分隔页
        if insert_dividers:
            slide = prs.slides.add_slide(blank_layout)
            _add_rect(slide, 0, 0, 13.333, 7.5, DEEP_BLUE)
            # 大数字(品牌蓝,140pt,放在偏上)
            _add_text(slide, 0.8, 1.5, 12, 2.5, f"{s_idx+1:02d}", size=140, bold=True, color=BRAND)
            # 装饰条
            _add_rect(slide, 0.8, 4.6, 1.0, 0.04, BRAND)
            # 章节名(白字,36pt)
            _add_text(slide, 0.8, 4.8, 12, 0.9, sec.get("name", ""), size=36, bold=True, color=WHITE)
            # 英文副名(辅蓝,14pt)
            if sec.get("name_en"):
                _add_text(slide, 0.8, 5.7, 12, 0.5, sec.get("name_en", ""), size=14, color=ACCENT)
            page_no += 1
            _add_footer(slide, page_no, total_pages)

        # 该章节下的内容页
        for sd in sec.get("slides", []):
            slide = prs.slides.add_slide(blank_layout)
            _add_rect(slide, 0, 0, 13.333, 0.18, BRAND)
            # 标题
            _add_text(slide, 0.6, 0.45, 12, 0.6, sd.get("title", ""), size=24, bold=True, color=DEEP_BLUE)
            # 引言(可选)
            if sd.get("intro"):
                _add_text(slide, 0.6, 1.15, 12, 0.5, sd.get("intro", ""), size=14, color=GRAY_TEXT)

            stype = sd.get("type", "bullets")

            # --------------- KPI 页 ---------------
            if stype == "kpi":
                kpis = sd.get("kpis", [])[:3]  # 最多 3 个
                kpi_w = 3.9
                gap = 0.25
                total_w = kpi_w * len(kpis) + gap * (len(kpis) - 1)
                start_x = (13.333 - total_w) / 2
                for i, k in enumerate(kpis):
                    x = start_x + i * (kpi_w + gap)
                    # 卡片底
                    _add_rect(slide, x, 2.1, kpi_w, 1.8, LIGHT_BG)
                    # 左侧蓝色边条
                    _add_rect(slide, x, 2.1, 0.04, 1.8, BRAND)
                    # 标签(12pt 灰)
                    _add_text(slide, x + 0.2, 2.25, kpi_w - 0.3, 0.4,
                              k.get("label", ""), size=12, color=GRAY_TEXT)
                    # 数值(28pt 加粗深蓝)
                    _add_text(slide, x + 0.2, 2.7, kpi_w - 0.3, 0.7,
                              k.get("value", ""), size=28, bold=True, color=DEEP_BLUE)
                    # 同比(11pt 灰)
                    if k.get("delta"):
                        _add_text(slide, x + 0.2, 3.4, kpi_w - 0.3, 0.4,
                                  k.get("delta", ""), size=11, color=BRAND)
                # 关键解读(bullets 形式)
                bullets = sd.get("bullets", [])
                if bullets:
                    cur_y = 4.3
                    for b in bullets:
                        # 自定义蓝色方块 bullet
                        _add_rect(slide, 0.7, cur_y + 0.12, 0.1, 0.1, BRAND)
                        _add_text(slide, 0.95, cur_y, 11.5, 0.5, b, size=14, color=DARK_TEXT)
                        cur_y += 0.55

            # --------------- 表格页 ---------------
            elif stype == "table":
                headers = sd.get("headers", [])
                rows = sd.get("rows", [])
                if headers and rows:
                    n_cols = len(headers)
                    table_x = 0.6
                    table_w = 12.1
                    col_w = table_w / n_cols
                    header_h = 0.5
                    row_h = 0.5
                    # 表头(深蓝底白字)
                    _add_rect(slide, table_x, 2.1, table_w, header_h, DEEP_BLUE)
                    for ci, htxt in enumerate(headers):
                        _add_text(slide, table_x + ci * col_w + 0.2, 2.18, col_w - 0.4, header_h - 0.1,
                                  htxt, size=13, bold=True, color=WHITE,
                                  align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT,
                                  anchor=MSO_ANCHOR.MIDDLE)
                    # 数据行
                    for ri, row in enumerate(rows[:12]):
                        ry = 2.1 + header_h + ri * row_h
                        # 偶数行浅蓝底
                        if ri % 2 == 1:
                            _add_rect(slide, table_x, ry, table_w, row_h, LIGHT_BG)
                        for ci in range(n_cols):
                            cell_text = row[ci] if ci < len(row) else ""
                            _add_text(slide, table_x + ci * col_w + 0.2, ry + 0.05,
                                      col_w - 0.4, row_h - 0.1,
                                      str(cell_text), size=13, color=DARK_TEXT,
                                      align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT,
                                      anchor=MSO_ANCHOR.MIDDLE)
                # 表格下解读
                bullets = sd.get("bullets", [])
                if bullets:
                    cur_y = 2.1 + header_h + min(len(rows), 12) * row_h + 0.4
                    for b in bullets:
                        _add_rect(slide, 0.7, cur_y + 0.12, 0.1, 0.1, BRAND)
                        _add_text(slide, 0.95, cur_y, 11.5, 0.5, b, size=14, color=DARK_TEXT)
                        cur_y += 0.55

            # --------------- 要点页(默认) ---------------
            else:
                bullets = sd.get("bullets", [])
                cur_y = 2.1
                for b in bullets:
                    _add_rect(slide, 0.7, cur_y + 0.12, 0.1, 0.1, BRAND)
                    _add_text(slide, 0.95, cur_y, 11.5, 0.5, b, size=15, color=DARK_TEXT)
                    cur_y += 0.6
                # 若有 intro,补一行小字
                if not bullets and sd.get("intro"):
                    _add_text(slide, 0.7, 2.1, 12, 0.6, sd["intro"], size=15, color=DARK_TEXT)

            page_no += 1
            _add_footer(slide, page_no, total_pages)

    # ------------------------------------------------------------------
    # 页型 6:结束页
    # ------------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.333, 7.5, DEEP_BLUE)
    # 装饰条
    _add_rect(slide, 0, 5.4, 13.333, 0.04, BRAND)
    # 谢谢聆听(54pt 白字)
    _add_text(slide, 0, 2.6, 13.333, 1.5, "谢  谢  聆  听", size=54, bold=True,
              color=WHITE, align=PP_ALIGN.CENTER)
    # Q & A
    _add_text(slide, 0, 4.0, 13.333, 0.6, "Q & A", size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    # 联系方式
    if closing_contact:
        _add_text(slide, 0, 5.8, 13.333, 0.5, closing_contact, size=12, color=ACCENT,
                  align=PP_ALIGN.CENTER)
    # 汇报主题(底部)
    _add_text(slide, 0, 6.8, 13.333, 0.4, title, size=11, color=ACCENT, align=PP_ALIGN.CENTER)
    # 结束页不加 _add_footer(深底,色条会让页码看不清)

    prs.save(path)
    return f"PPT 已生成:{path}(共 {total_pages} 页,含封面/目录/{len(sections)} 个章节/{'分隔页+' if insert_dividers else ''}结束页)"


@tool
def load_skill(name: str) -> str:
    """加载指定 skill 的完整内容。
    name 取自系统提示中"可用技能"列表,如 'excel-beautification' 或 'ppt-generation'。
    返回该 skill 的 SKILL.md 正文(去掉 YAML frontmatter)。
    """
    return SKILL_LOADER.get_content(name)


# ----------------------------------------------------------------------
# 3. 专家智能体:总结/翻译是纯 LLM(无工具),其余各挂一个工具
# ----------------------------------------------------------------------
summarizer = create_agent(
    model,
    tools=[],
    system_prompt="你是总结专家。把用户提供的内容提炼成简洁、结构清晰的中文总结,抓住要点。",
    name="summarizer",
)

translator = create_agent(
    model,
    tools=[],
    system_prompt="你是翻译专家。根据用户要求在中英文(或指定语言)之间准确翻译,保持原意、语气和专有名词。",
    name="translator",
)

emailer = create_agent(
    model,
    tools=[send_email, request_human_input],
    system_prompt=(
        "你负责发送邮件。流程固定为:\n"
        "1. 根据用户意图整理出收件人、主题和正文;\n"
        "2. **先调用 request_human_input** 向用户确认(系统会自动展示邮件预览);\n"
        "3. 用户选择 Approve → 调用 send_email 实际发送;\n"
        "4. 用户选择 Reject(或其它) → 不发送,礼貌说明并询问如何调整。"
    ),
    name="emailer",
)

ppt_maker = create_agent(
    model,
    tools=[make_ppt, make_business_ppt, load_skill],
    system_prompt = f"""
    你是一名专业的 PPT 制作助手。

    任务：
    根据用户提供的内容，先规划 PPT 结构，再调用合适的工具生成演示文稿。

    工具选择规则（务必遵循）：
    1. **商务汇报 / 季度复盘 / 销售提案 / 项目汇报 / 产品介绍 / 培训分享 等结构化演示** →
       必须使用 `make_business_ppt`。该工具按 16:9 + 品牌蓝 + 6 种页型(封面/目录/章节分隔/KPI/表格/结束)
       的统一规范生成,视觉专业、可直接交付。
    2. **非结构化 / 自由排版 / 海报 / 邀请函 等单页设计** →
       使用 `make_ppt`(简单封面+标题/内容两栏)。
    3. 不确定时,默认走 `make_business_ppt`。

    工作目录：
    {WORKDIR}

    可用技能：
    {SKILL_LOADER.get_descriptions()}

    使用 `make_business_ppt` 时的输入约定:
    - sections 是 3–6 个章节,每个章节内含 1–3 张内容页
    - content page type 必须是 `kpi` / `bullets` / `table` 之一
    - KPI 页 kpis 字段最多 3 个(超出将自动截断)
    - 表格页最多 12 行(超出将自动截断)
    - 字段顺序: filename → title → subtitle(可选) → presenter(可选) → date(可选) → sections → closing_contact(可选)
    - 调用前必须先用一段文字向用户说明 PPT 章节结构,确认后一次性调用工具

    输出要求：
    - 先生成完整的 PPT 大纲(章节 + 每页类型)。
    - 确认结构后调用对应工具。
    - 最终 PPT 应专业、简洁、易于展示。
    """,
    name="ppt_maker",
)

excel_maker = create_agent(
    model,
    tools=[make_excel, load_skill],
    system_prompt = f"""
    你是一名专业的 Excel 表格生成助手。

    任务：
    根据用户提供的内容，先整理表格结构，再调用 `make_excel` 生成 Excel 文件。

    要求：
    1. 分析用户需求，提取结构化数据。
    2. 将数据整理为：
    - 表头（headers）
    - 数据行（rows）
    3. 表头命名清晰、规范、易于理解。
    4. 数据格式保持一致，确保行列对应正确。
    5. 对于非结构化内容，先进行归纳整理，再转换为表格形式。
    6. 对于统计、分析或业务数据，合理设计字段，提升可读性和后续分析价值。
    7. 在处理不熟悉的领域前，优先使用 `load_skill` 获取相关知识。

    工作目录：
    {WORKDIR}

    可用技能：
    {SKILL_LOADER.get_descriptions()}

    输出要求：
    - 先生成完整的表格结构（headers 和 rows）。
    - 检查数据完整性和一致性。
    - 然后调用 `make_excel` 生成 Excel 文件。
    - 最终表格应规范、准确、易于查看和后续处理。
    """,
    name="excel_maker",
)


# ----------------------------------------------------------------------
# 4. 主管:负责把用户任务路由给合适的专家,支持多步串联
#    注意:supervisor 编排这一层对库版本较敏感,若组合报错见文件末尾说明
# ----------------------------------------------------------------------
supervisor = create_supervisor(
    agents=[summarizer, translator, emailer, ppt_maker, excel_maker],
    model=model,
    # add_handoff_back_messages=False  → don't pollute the conversation
    # history with "Transferring back to supervisor" / "Successfully
    # transferred back to supervisor" AIMessage+ToolMessage pairs every
    # time a sub-agent returns. They're internal plumbing.
    add_handoff_back_messages=False,
    # pre_model_hook: structural guard against the supervisor LLM looping
    # back to the same agent. After the first transfer_to_* call we collapse
    # the visible history so the LLM can only end or pick a *different*
    # agent — never re-invoke the same one.
    pre_model_hook=_supervisor_pre_model_hook,
    prompt=(
        "你是团队主管,负责把用户任务分派给合适的专家,自己不直接做事:\n"
        "- 内容总结/提炼 → summarizer\n"
        "- 语言翻译 → translator\n"
        "- 发送邮件 → emailer\n"
        "- 制作 PPT → ppt_maker\n"
        "- 制作 Excel → excel_maker\n"
        "\n"
        "硬性规则(务必遵守):\n"
        "1. **每个专家对每个子任务只调用一次**——不要重复把同一个专家叫回来追问或'再确认'。\n"
        "   一旦专家返回结果,直接基于该结果决定下一步或汇报给用户。\n"
        "2. **多步任务按需串行**:例如'先总结再做PPT'就是 summarizer 一次 → ppt_maker 一次,中间不再回头。\n"
        "3. **完成后立即收尾**:所有需要的专家都跑完后,直接向用户汇报最终结果,不要再触发任何 transfer。\n"
        "4. 如果专家的结果已经满足用户需求,不要再开新轮,直接给最终回复即可。\n"
    ),
).compile()


# ----------------------------------------------------------------------
# 5. 运行示例
# ----------------------------------------------------------------------
def run(user_input: str):
    result = supervisor.invoke({"messages": [{"role": "user", "content": user_input}]})
    # 打印最终回复
    print(result["messages"][-1].content)
    return result


# ----------------------------------------------------------------------
# 6. 流式生成器:同时输出文本 token + agent 切换 + tool 调用事件
# ----------------------------------------------------------------------
async def astream_supervisor(user_input: str):
    """异步生成器,yield dict 形式的流式事件。

    事件类型:
      {"event": "agent",   "node": "summarizer"}                              # agent 开始/切换
      {"event": "token",   "node": "summarizer", "content": "..."}            # 文本 token
      {"event": "tool_call",  "node": "ppt_maker", "tool": "load_skill", "args": {"name": "ppt-generation"}}
      {"event": "tool_result","node": "ppt_maker", "tool": "load_skill", "preview": "..."}
      {"event": "human_input_required", "request_id": "...", "question": "...", "options": [...]}  # HITL 暂停
      {"event": "done"}                                                       # 流结束
      {"event": "error",   "message": "..."}                                  # 异常
    """
    inputs = {"messages": [{"role": "user", "content": user_input}]}
    seen_nodes: set[str] = set()
    queue: asyncio.Queue = asyncio.Queue()
    loop = asyncio.get_running_loop()
    _hitl_broker.attach(queue, loop)

    async def _supervisor_to_queue() -> None:
        """把 supervisor 的 stream 事件转成统一 dict,放进共享 queue。"""
        try:
            async for mode, data in supervisor.astream(
                inputs, stream_mode=["messages", "updates"]
            ):
                if mode == "messages":
                    chunk, meta = data
                    node = meta.get("langgraph_node", "unknown")
                    if node not in seen_nodes:
                        seen_nodes.add(node)
                        await queue.put({"event": "agent", "node": node})
                    content = ""
                    if hasattr(chunk, "content"):
                        c = chunk.content
                        content = c if isinstance(c, str) else "".join(
                            p.get("text", "") if isinstance(p, dict) else str(p) for p in c
                        )
                    if content:
                        await queue.put({"event": "token", "node": node, "content": content})

                elif mode == "updates":
                    for node_name, state in data.items():
                        msgs = state.get("messages", []) if isinstance(state, dict) else []
                        for msg in msgs:
                            if hasattr(msg, "tool_calls") and msg.tool_calls:
                                for tc in msg.tool_calls:
                                    tool_name = tc.get("name", "unknown")
                                    # LangGraph's supervisor injects internal
                                    # handoff tools (transfer_to_X / transfer_back_to_X).
                                    # Those are plumbing the user never asked about,
                                    # so we skip emitting them as tool_call events.
                                    if _is_handoff_tool(tool_name):
                                        continue
                                    await queue.put({
                                        "event": "tool_call",
                                        "node": node_name,
                                        "tool": tool_name,
                                        "args": tc.get("args", {}),
                                    })
                            if msg.__class__.__name__ == "ToolMessage":
                                tool_name = getattr(msg, "name", "unknown")
                                if _is_handoff_tool(tool_name):
                                    # Same handoff-tool skip on the result side.
                                    continue
                                preview = str(getattr(msg, "content", ""))[:200]
                                await queue.put({
                                    "event": "tool_result",
                                    "node": node_name,
                                    "tool": tool_name,
                                    "preview": preview,
                                })
            await queue.put({"event": "done"})
        except Exception as e:
            await queue.put({"event": "error", "message": str(e)})

    sup_task = asyncio.create_task(_supervisor_to_queue())
    try:
        while True:
            evt = await queue.get()
            yield evt
            if evt.get("event") in ("done", "error"):
                break
    finally:
        _hitl_broker.detach()
        if not sup_task.done():
            sup_task.cancel()
            try:
                await sup_task
            except (asyncio.CancelledError, Exception):
                pass


if __name__ == "__main__":
    # 示例 1:纯总结
    run("帮我把下面这段话总结成三点:人工智能正在改变各行各业……")

    # 示例 2:多步 —— 总结后做成 PPT
    # run("把这篇文章总结一下,然后做成一个3页的PPT:……")

    # 示例 3:做 Excel
    # run("帮我做一个销售表,列是 月份、销售额,数据:1月100,2月200,3月150")


# ======================================================================
# 排错说明
# ----------------------------------------------------------------------
# 1) 若 `from langchain.agents import create_agent` 报错:说明 langchain 版本偏旧,
#    升级:pip install -U langchain langgraph
# 2) 若 create_supervisor 与 create_agent 组合时报参数/类型错误(这层 API 版本敏感):
#    - 备选 A:把 create_agent 换成 langgraph.prebuilt.create_react_agent(会有弃用警告但能跑)
#    - 备选 B:不用 langgraph-supervisor 库,改用 StateGraph + 手写 handoff 工具自建主管
#    把你的报错和 `pip show langchain langgraph langgraph-supervisor` 版本发我,我帮你对症改。
# ======================================================================